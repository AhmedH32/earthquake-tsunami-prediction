from pptx import Presentation


def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Earthquake-Induced Tsunami Prediction"
    slide.placeholders[1].text = "Spatial Engineering & Cost-Asymmetric Optimization\nAhmed Hassan | ML Capstone"

    # Slide 2: The Data
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Data & The Danger"
    tf = slide.placeholders[1].text_frame
    tf.text = "Curated dataset: 764 oceanic seismic events (61% Normal vs 39% Tsunami)"
    tf.add_paragraph().text = "Problem: Software defaults to 50/50 thresholds and symmetrical metrics"
    tf.add_paragraph().text = "Solution: Optimized strictly for the F2-score (penalizes missed tsunamis 4x harder)"

    # Slide 3: Engineering
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Spatial-Temporal Engineering"
    tf = slide.placeholders[1].text_frame
    tf.text = "Raw Lat/Lon causes 'Anti-Meridian Distortion' in the Pacific"
    tf.add_paragraph().text = "Fix: Projected to 3D Cartesian Manifold (x, y, z)"
    tf.add_paragraph().text = "Fix: Cyclical encoding for the 'month' feature"
    tf.add_paragraph().text = "Proof: Feature ablation boosted baseline F2 from 0.6401 to 0.7288"

    # Slide 4: Paradox
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Architectures & The Paradox"
    tf = slide.placeholders[1].text_frame
    tf.text = "KNN: Baseline density estimator (Failed due to tabular complexity)"
    tf.add_paragraph().text = "Decision Tree: Interpretable rules (Won baseline F2: 0.8900)"
    tf.add_paragraph().text = "Random Forest: High-capacity manifold learner"
    tf.add_paragraph().text = "The Paradox: RF achieved 0.9497 AUC, but ensemble averaging squashed probabilities below 0.5"

    # Slide 5: Visual Proof
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Visual Proof (ROC Curve)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Random Forest perfectly separates the physics of tsunamis."
    tf.add_paragraph().text = "[ Drag and Drop combined_roc_curve.png here ]"

    # Slide 6: Calibration
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Operational Calibration"
    tf = slide.placeholders[1].text_frame
    tf.text = "Decoupled structural training from operational deployment"
    tf.add_paragraph().text = "Executed leak-free coordinate sweep on training data"
    tf.add_paragraph().text = "Found optimal safety threshold: tau* = 0.3236"
    tf.add_paragraph().text = "Result: Test F2-score skyrocketed to 0.9223 (Caught 57 of 59 tsunamis)"

    # Slide 7: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusion & Future Work"
    tf = slide.placeholders[1].text_frame
    tf.text = "Domain-specific geometric engineering is mandatory for geographic data"
    tf.add_paragraph().text = "Decoupling capacity from calibration unlocks ensemble models"
    tf.add_paragraph().text = "Future Work: XGBoost (custom loss) and LSTMs (time-series clustering)"

    prs.save('Capstone_Presentation.pptx')
    print("Presentation created successfully!")

if __name__ == "__main__":
    create_presentation()